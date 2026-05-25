from fastapi import FastAPI, Query, HTTPException, Body, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from contextlib import asynccontextmanager
from apscheduler.schedulers.background import BackgroundScheduler
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
# --- PPTX IMPORTS ---
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
# --- DOCX IMPORTS ---
from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
# --- EXCEL IMPORTS ---
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
# --- IMAGE & AI IMPORTS ---
from PIL import Image, ImageOps
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from typing import List, Optional
import subprocess
import os
import json
import sys
import datetime
import glob
import pandas as pd
import io
import requests
import learner
import re
import urllib3
import pickle
import numpy as np
import threading
import secrets
import platform
import hashlib
from pathlib import Path
from sentence_transformers import SentenceTransformer
from fastapi import BackgroundTasks
from concurrent.futures import ThreadPoolExecutor
from threading import Semaphore

# ==========================================
# AI GATEKEEPER MODEL LOADING
# ==========================================
print("Waking up the AI Gatekeeper...")

BOUNCER_MODEL_FILENAMES = {
    "default": "bouncer_model.pkl",
    "broadcast": "bouncer_model_broadcast.pkl",
}

try:
    current_dir = os.path.dirname(os.path.abspath(__file__))
    model_folder = os.path.join(current_dir, "local_miniLM_model")
    bouncer_embedder = SentenceTransformer(model_folder)
    bouncer_models = {}

    for profile_name, model_filename in BOUNCER_MODEL_FILENAMES.items():
        model_path = os.path.join(current_dir, model_filename)
        if not os.path.exists(model_path):
            print(f"No {profile_name} bouncer found yet: {model_filename}")
            continue
        try:
            with open(model_path, "rb") as f:
                bouncer_models[profile_name] = pickle.load(f)
            print(f"Loaded {profile_name} bouncer: {model_filename}")
        except Exception as model_error:
            print(
                f"Could not load {profile_name} bouncer "
                f"({model_filename}): {model_error}"
            )

    bouncer_model = bouncer_models.get("default")
    if bouncer_models:
        print("AI Gatekeeper is awake and profile-aware.")
    else:
        print("No bouncer models loaded. Scanning without filter.")
except Exception as e:
    print(f"Warning: Gatekeeper not found. Scanning without filter. Error: {e}")
    bouncer_embedder = None
    bouncer_models = {}
    bouncer_model = None

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# ==========================================
# CONFIGURATION
# ==========================================
MORNING_KEYWORDS = (
    "OpenAI , Robot , Samsung , LG , Sony , Nvidia , TCL , OLED , QNED , "
    "Artificial Intelligence, chatGPT , Anthropic , Claude , Gemini , LED , "
    "Robotics , Television , TV , display , Grok , GPU , Processor , Jio , TPU"
)
BROADCAST_MORNING_KEYWORDS = (
    "DTH, Cable TV, IPTV, Broadcast, Digital terrestrial transmission, DTT, "
    "DVB S, DVB S2, DVB C, DVB C2, DVB T, DVB T2, conditional access system, "
    "digital rights management, FAST, OTT, Connected TV, Tuner, Set top box, "
    "Linear ad insertion, Linear ads, TRAI, MIB, broadcast regulation, HBB TV, "
    "DVB I, 5G broadcast, D2M"
)
DIRECTOR_KEY = os.environ.get("DIRECTOR_KEY", "1357")
ANALYTICS_KEY = os.environ.get("ANALYTICS_KEY", DIRECTOR_KEY)
ROOT_DIR = os.getcwd()
HISTORY_DIR = os.path.join(ROOT_DIR, "history_archive")
MANUAL_LOG_FILE = os.path.join(ROOT_DIR, "manual_search_logs.xlsx")
WORKFLOW_FILE = os.path.join(ROOT_DIR, "workflow_store.json")
TRAINING_FILE = os.path.join(ROOT_DIR, "trainingData.json")
NOT_INTERESTED_FILE = os.path.join(ROOT_DIR, "not_interested_store.json")
NOT_INTERESTED_EXPIRY_HOURS = 22
USAGE_TRACKER_FILE = os.path.join(ROOT_DIR, "usage_tracker.json")

# ==========================================
# PROFILE ROUTING AND STORAGE
# ==========================================
DEFAULT_PROFILE = "default"
BROADCAST_PROFILE = "broadcast"
BROADCAST_SPECIAL_IPS = {
    "107.109.202.212",
    "107.109.202.33",
}
ANALYTICS_ALLOWED_IPS = {
    ip.strip()
    for ip in os.environ.get(
        "ANALYTICS_ALLOWED_IPS",
        "127.0.0.1,::1,107.109.201.245",
    ).split(",")
    if ip.strip()
}

DEFAULT_SITES_FILE = os.path.join(ROOT_DIR, "sites.json")
BROADCAST_SITES_FILE = os.path.join(ROOT_DIR, "sites_broadcast.json")
INTELLIGENCE_STORE_DIR = os.path.join(ROOT_DIR, "intelligence_store")

PROFILE_CONFIGS = {
    DEFAULT_PROFILE: {
        "label": "Default Intelligence",
        "keywords": MORNING_KEYWORDS,
        "sites_file": DEFAULT_SITES_FILE,
        "history_dir": os.path.join(INTELLIGENCE_STORE_DIR, DEFAULT_PROFILE, "history"),
        "use_bouncer": True,
    },
    BROADCAST_PROFILE: {
        "label": "Broadcast Intelligence",
        "keywords": BROADCAST_MORNING_KEYWORDS,
        "sites_file": BROADCAST_SITES_FILE,
        "history_dir": os.path.join(INTELLIGENCE_STORE_DIR, BROADCAST_PROFILE, "history"),
        "use_bouncer": True,
    },
}
WORKFLOW_FILES = {
    DEFAULT_PROFILE: WORKFLOW_FILE,
    BROADCAST_PROFILE: os.path.join(ROOT_DIR, "workflow_store_broadcast.json"),
}
NOT_INTERESTED_FILES = {
    DEFAULT_PROFILE: NOT_INTERESTED_FILE,
    BROADCAST_PROFILE: os.path.join(ROOT_DIR, "not_interested_store_broadcast.json"),
}
TRAINING_FILES = {
    DEFAULT_PROFILE: TRAINING_FILE,
    BROADCAST_PROFILE: os.path.join(ROOT_DIR, "trainingData_broadcast.json"),
}
BOUNCER_MODEL_FILES = {
    DEFAULT_PROFILE: os.path.join(ROOT_DIR, "bouncer_model.pkl"),
    BROADCAST_PROFILE: os.path.join(ROOT_DIR, "bouncer_model_broadcast.pkl"),
}
REGION_LEARNING_FILES = {
    DEFAULT_PROFILE: os.path.join(ROOT_DIR, "region_learning.json"),
    BROADCAST_PROFILE: os.path.join(ROOT_DIR, "region_learning_broadcast.json"),
}

TEAM_IP_MAP = {
    "107.109.202.151": "Shreya Gupta",
    "107.109.201.245": "Vineet Singh",
    "107.109.202.48": "Shivani Goyal",
    "107.109.201.45": "ASHOK JAIN",
    "107.109.201.66": "Priya Arora",
    "107.109.201.83": "Vinod Sati",
    "107.109.201.58": "Ravi Kant Bansal",
    "107.109.202.178": "Traxon PC(HOST)",
    "107.109.201.101": "Utkarsh Tiwari",
}


def get_client_ip(request: Request = None):
    if not request:
        return "unknown"
    forwarded = request.headers.get("x-forwarded-for")
    if forwarded:
        return forwarded.split(",")[0].strip()
    real_ip = request.headers.get("x-real-ip")
    if real_ip:
        return real_ip.strip()
    return request.client.host if request.client else "unknown"


def get_profile_for_request(request: Request = None):
    if request:
        requested_profile = (
            request.headers.get("x-sense-profile")
            or request.query_params.get("profile")
            or ""
        ).strip().lower()
        if requested_profile in PROFILE_CONFIGS:
            return requested_profile

    return (
        BROADCAST_PROFILE
        if get_client_ip(request) in BROADCAST_SPECIAL_IPS
        else DEFAULT_PROFILE
    )


def get_active_profile_name(request: Request = None):
    return get_profile_for_request(request)


def get_profile_config(profile: str):
    return PROFILE_CONFIGS.get(profile, PROFILE_CONFIGS[DEFAULT_PROFILE])


def get_sites_file_for_profile(profile: str):
    return get_profile_config(profile)["sites_file"]


def get_profile_history_dir(profile: str):
    history_dir = get_profile_config(profile)["history_dir"]
    os.makedirs(history_dir, exist_ok=True)
    return history_dir


def ensure_profile_storage():
    for profile in PROFILE_CONFIGS:
        os.makedirs(get_profile_history_dir(profile), exist_ok=True)


def get_profile_history_files(profile: str, include_legacy_default: bool = True):
    files = glob.glob(os.path.join(get_profile_history_dir(profile), "*.json"))
    if include_legacy_default and profile == DEFAULT_PROFILE:
        files.extend(glob.glob(os.path.join(HISTORY_DIR, "*.json")))
    return list(dict.fromkeys(os.path.abspath(file_path) for file_path in files))


def resolve_profile_history_file(filename: str, profile: str):
    safe_name = Path(filename).name
    profile_path = os.path.join(get_profile_history_dir(profile), safe_name)
    if os.path.exists(profile_path):
        return profile_path
    legacy_path = os.path.join(HISTORY_DIR, safe_name)
    if profile == DEFAULT_PROFILE and os.path.exists(legacy_path):
        return legacy_path
    return None


def get_latest_briefing_file_for_profile(profile: str):
    files = glob.glob(os.path.join(get_profile_history_dir(profile), "briefing_*.json"))
    if profile == DEFAULT_PROFILE and not files:
        files = glob.glob(os.path.join(HISTORY_DIR, "briefing_*.json"))
    valid_files = []
    for file_path in files:
        try:
            with open(file_path, "r", encoding="utf-8") as file_obj:
                if json.load(file_obj):
                    valid_files.append(file_path)
        except Exception:
            continue
    return max(valid_files, key=os.path.getmtime) if valid_files else None


def get_workflow_file_for_request(request: Request = None):
    return WORKFLOW_FILES.get(get_profile_for_request(request), WORKFLOW_FILE)


def get_not_interested_file_for_profile(profile: str):
    return NOT_INTERESTED_FILES.get(profile, NOT_INTERESTED_FILE)


def get_training_file_for_profile(profile: str):
    return TRAINING_FILES.get(profile, TRAINING_FILE)


def get_bouncer_model_file_for_profile(profile: str):
    return BOUNCER_MODEL_FILES.get(profile, BOUNCER_MODEL_FILES[DEFAULT_PROFILE])


def get_bouncer_model_for_profile(profile: str):
    return bouncer_models.get(profile) if bouncer_embedder is not None else None


def get_team_owner_for_ip(ip: str):
    return TEAM_IP_MAP.get(str(ip or "").strip())


def is_analytics_allowed_ip(ip: str) -> bool:
    return str(ip or "").strip() in ANALYTICS_ALLOWED_IPS


def require_analytics_access(request: Request, key: str = None):
    ip = get_client_ip(request)

    if not is_analytics_allowed_ip(ip):
        raise HTTPException(status_code=403, detail="Analytics is not enabled for this network.")

    if key != ANALYTICS_KEY:
        raise HTTPException(status_code=403, detail="Invalid analytics key.")

    return ip


def get_profile_debug_info(profile: str):
    config = get_profile_config(profile)
    latest_file = get_latest_briefing_file_for_profile(profile)
    return {
        "profile": profile,
        "label": config["label"],
        "sites_file": config["sites_file"],
        "sites_file_exists": os.path.exists(config["sites_file"]),
        "history_dir": config["history_dir"],
        "history_dir_exists": os.path.exists(config["history_dir"]),
        "latest_briefing": os.path.basename(latest_file) if latest_file else None,
        "bouncer_model_file": get_bouncer_model_file_for_profile(profile),
        "bouncer_model_exists": os.path.exists(get_bouncer_model_file_for_profile(profile)),
        "training_file": get_training_file_for_profile(profile),
        "training_file_exists": os.path.exists(get_training_file_for_profile(profile)),
        "not_interested_file": get_not_interested_file_for_profile(profile),
        "not_interested_file_exists": os.path.exists(get_not_interested_file_for_profile(profile)),
    }


BOUNCER_LOW_PRIORITY_THRESHOLD = 0.45
BOUNCER_HARD_DROP_THRESHOLD = 0.60

# ==========================================
# --- THREAD POOL FOR ML INFERENCE ---
# ==========================================
ml_executor = ThreadPoolExecutor(max_workers=2, thread_name_prefix="ml_worker")

# ==========================================
# --- MULTI-USER JOB TRACKING ---
# ==========================================
crawl_semaphore = Semaphore(3)
active_jobs = {}

SCHEDULER_STATUS = {
    "is_active": False,
    "message": "System Ready.",
    "mode": "idle",
}

scheduler_lock = threading.Lock()
file_lock = threading.Lock()
train_lock = threading.Lock()
not_interested_lock = threading.Lock()
tracker_lock = threading.Lock()
region_learning_lock = threading.Lock()
dropped_lock = threading.Lock()
opinion_lock = threading.Lock()
insight_cache_lock = threading.Lock()
insight_cache = {}

CLOSE_FDS = platform.system() != "Windows"

if not os.path.exists(HISTORY_DIR):
    os.makedirs(HISTORY_DIR)

# ==========================================
# --- INITIALIZE HUGGING FACE OPINION ENGINE ---
# ==========================================
print("Initializing Hugging Face Opinion Engine (Local)...")
try:
    model_path = "./flan-t5-local"
    tokenizer = AutoTokenizer.from_pretrained(model_path, local_files_only=True)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_path, local_files_only=True)
    print("Opinion Engine Ready.")
except Exception as e:
    print(f"Opinion Engine Failed to Load: {e}")
    tokenizer = None
    model = None


def generate_opinion(text):
    if not model or not tokenizer or not text:
        return "Insight generation unavailable."
    try:
        prompt = f"Briefly analyze this news and give a one-sentence professional opinion: {text[:500]}"
        with opinion_lock:
            inputs = tokenizer(
                prompt, return_tensors="pt", max_length=512, truncation=True
            )
            outputs = model.generate(**inputs, max_new_tokens=40, do_sample=False)
        return tokenizer.decode(outputs[0], skip_special_tokens=True)
    except Exception as e:
        print(f"Insight Gen Error: {e}")
        return "Could not compute insight."


def fallback_why_it_matters(item):
    source_count = int(item.get("source_count", 1) or 1)
    category = str(item.get("category", "technology intelligence") or "technology intelligence")
    category_lower = category.lower()
    if "broadcast" in category_lower:
        consequence = "It could reshape distribution reach, rights economics, and regulatory priorities for broadcast operators."
    elif "ai" in category_lower:
        consequence = "It may shift product capability, compute demand, and competitive positioning for AI teams."
    elif any(term in category_lower for term in ["display", "television", "device"]):
        consequence = "It may influence product roadmaps, supplier choices, and near-term competitive differentiation."
    else:
        consequence = "It may change competitive priorities, investment choices, or execution risk for decision-makers."
    return (
        f"{consequence} The signal is supported by {source_count} "
        f"source{'s' if source_count != 1 else ''}."
    )


def is_weak_generated_insight(insight, title):
    generated_words = re.findall(r"[a-z0-9]+", str(insight or "").lower())
    title_words = set(re.findall(r"[a-z0-9]+", str(title or "").lower()))
    if len(generated_words) < 9:
        return True
    overlap = sum(1 for word in generated_words if word in title_words)
    return overlap / max(len(generated_words), 1) > 0.7


def generate_why_it_matters(item, profile=DEFAULT_PROFILE):
    title = str(item.get("title", "") or "").strip()
    summary = str(
        item.get("master_summary")
        or item.get("summary")
        or item.get("snippet")
        or ""
    ).strip()
    cache_key = hashlib.sha256(
        f"{profile}|{title}|{summary[:1000]}".encode("utf-8")
    ).hexdigest()

    with insight_cache_lock:
        cached = insight_cache.get(cache_key)
    if cached:
        return cached[0], f"{cached[1]}-cache"

    if not model or not tokenizer or not summary:
        insight = fallback_why_it_matters(item)
        source = "fallback"
    else:
        prompt = (
            "You are preparing an executive intelligence briefing. "
            "Complete the sentence 'This matters because...' in one concise sentence, "
            "describing strategic impact, risk, opportunity, or market consequence, "
            "not repeating the headline. "
            f"Title: {title}. Summary: {summary[:800]}"
        )
        try:
            with opinion_lock:
                inputs = tokenizer(
                    prompt,
                    return_tensors="pt",
                    max_length=512,
                    truncation=True,
                )
                outputs = model.generate(
                    **inputs,
                    max_new_tokens=64,
                    do_sample=False,
                )
            insight = tokenizer.decode(outputs[0], skip_special_tokens=True).strip()
            source = "flan-t5-local"
            if is_weak_generated_insight(insight, title):
                insight = fallback_why_it_matters(item)
                source = "fallback-after-flan"
        except Exception as e:
            print(f"Why It Matters Gen Error: {e}")
            insight = fallback_why_it_matters(item)
            source = "fallback"

    with insight_cache_lock:
        if len(insight_cache) >= 1000:
            insight_cache.pop(next(iter(insight_cache)))
        insight_cache[cache_key] = (insight, source)
    return insight, source


# ==========================================
# --- TEAM ROUTING LOGIC ---
# ==========================================
def determine_target_team(title, summary):
    text = (title + " " + summary).lower()
    if any(k in text for k in ["cloud", "server", "aws", "azure"]):
        return "Cloud Team"
    elif any(k in text for k in ["hardware", "chip", "semiconductor", "nvidia"]):
        return "Hardware Team"
    elif any(k in text for k in ["robot", "robotics", "automation"]):
        return "Robotics Team"
    elif any(k in text for k in ["tv", "oled", "display", "tizen", "streaming"]):
        return "TV & Display Team"
    elif any(
        k in text
        for k in ["broadcast", "dth", "iptv", "ott", "dvb", "set top box", "trai", "mib"]
    ):
        return "Broadcast Team"
    else:
        return "ALL"


# ==========================================
# --- NEWS CATEGORIZATION ENGINE ---
# ==========================================
CATEGORY_MATRIX = {
    "AI Models": [
        "llm", "gpt", "gemini", "claude", "llama", "foundation model",
        "parameters", "openai", "anthropic", "neural network",
    ],
    "AI Agents": [
        "agent", "autonomous agent", "copilot", "ai assistant",
        "digital assistant", "virtual assistant",
    ],
    "Smart Features": [
        "smart feature", "intelligent feature", "auto-",
        "smart tracking", "adaptive", "predictive",
    ],
    "Form Factor": [
        "form factor", "foldable", "rollable", "wearable", "design",
        "chassis", "slimmer", "hinge",
    ],
    "New Product": [
        "launch", "unveil", "release", "debut", "announced",
        "new lineup", "introducing",
    ],
    "Robotics": [
        "robot", "humanoid", "boston dynamics", "automation",
        "bipedal", "robotic", "drone",
    ],
    "Services": [
        "subscription", "service", "cloud service",
        "platform as a service", "saas", "streaming",
    ],
    "Security": [
        "security", "privacy", "cybersecurity", "hack", "breach",
        "encryption", "knox", "malware",
    ],
    "Smart Home": [
        "smart home", "iot", "thermostat", "fridge", "appliance",
        "smartthings", "matter",
    ],
    "Display Tech": [
        "oled", "microled", "display", "screen", "monitor", "tv",
        "resolution", "nits", "panel",
    ],
    "Partnership": [
        "partnership", "collaboration", "team up", "teaming up",
        "joint venture", "partnered",
    ],
    "Research": [
        "research", "study", "paper", "breakthrough", "scientists",
        "developed a new", "laboratory",
    ],
    "Patent": [
        "patent", "trademark", "intellectual property", "uspto",
        "filed a patent",
    ],
    "Broadcasting": [
        "broadcast", "dth", "cable tv", "iptv", "dvb", "ott", "fast",
        "connected tv", "set top box", "tuner", "linear ad insertion",
        "trai", "mib", "broadcast regulation", "hbbtv", "5g broadcast", "d2m",
        "digital terrestrial transmission", "conditional access",
        "digital rights management",
    ],
    "AI Features": [
        "ai-powered", "generative ai", "genai", "ai capability",
        "ai tool", "ai update",
    ],
}


def assign_category(title, summary):
    text = (str(title) + " " + str(summary)).lower()
    scores = {cat: 0 for cat in CATEGORY_MATRIX}
    for category, keywords in CATEGORY_MATRIX.items():
        for kw in keywords:
            if kw in text:
                scores[category] += text.count(kw)
    best_category = max(scores, key=scores.get)
    if scores[best_category] == 0:
        return "Tech News"
    return best_category


# ==========================================
# --- USER-TAUGHT REGION CLASSIFICATION ---
# ==========================================
def normalize_region_label(value):
    candidate = str(value or "").strip().lower()
    if candidate == "local":
        return "Local"
    if candidate == "global":
        return "Global"
    return None


def normalize_region_keywords(value):
    if isinstance(value, list):
        raw_values = value
    else:
        raw_values = re.split(r"[,;\n]", str(value or ""))
    keywords = []
    for raw in raw_values:
        keyword = re.sub(r"\s+", " ", str(raw).strip().lower())
        if len(keyword) >= 2 and keyword not in keywords:
            keywords.append(keyword)
    return keywords[:20]


def get_region_learning_file_for_profile(profile: str):
    return REGION_LEARNING_FILES.get(profile, REGION_LEARNING_FILES[DEFAULT_PROFILE])


def load_region_learning(profile=DEFAULT_PROFILE):
    blank = {"Local": [], "Global": [], "corrections": []}
    learning_file = get_region_learning_file_for_profile(profile)
    if not os.path.exists(learning_file):
        return blank
    try:
        with open(learning_file, "r", encoding="utf-8") as f:
            stored = json.load(f)
        for region in ("Local", "Global"):
            blank[region] = normalize_region_keywords(stored.get(region, []))
        blank["corrections"] = stored.get("corrections", [])[-500:]
    except (json.JSONDecodeError, OSError, TypeError):
        pass
    return blank


def save_region_learning(data, profile=DEFAULT_PROFILE):
    with open(get_region_learning_file_for_profile(profile), "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)


def region_text_for_article(item):
    fields = [
        item.get("title", ""),
        item.get("master_summary", ""),
        item.get("summary", ""),
        item.get("snippet", ""),
        item.get("full_contents", ""),
        item.get("source", ""),
        " ".join(item.get("keywords_found", []) or []),
    ]
    return " ".join(str(field) for field in fields if field).lower()


def apply_learned_region(item, profile=DEFAULT_PROFILE):
    if not isinstance(item, dict):
        return item
    learned = load_region_learning(profile)
    title = str(item.get("title", "")).strip().lower()
    for correction in reversed(learned["corrections"]):
        if title and title == str(correction.get("title", "")).strip().lower():
            next_item = dict(item)
            next_item["region"] = correction["region"]
            next_item["region_basis"] = "User corrected"
            return next_item

    text = region_text_for_article(item)
    for region in ("Local", "Global"):
        matches = [keyword for keyword in learned[region] if keyword in text]
        if matches:
            next_item = dict(item)
            next_item["region"] = region
            next_item["region_basis"] = f"Learned keyword: {matches[0]}"
            return next_item
    return item


def apply_learned_regions(items, profile=DEFAULT_PROFILE):
    return [apply_learned_region(item, profile) for item in (items or [])]


# ==========================================
# --- HELPER: Robust Image Downloader ---
# ==========================================
def download_image_for_export(url, add_border=False):
    if not url:
        return None
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(
            url, headers=headers, timeout=(5, 10), verify=False
        )
        if response.status_code == 200:
            img = Image.open(io.BytesIO(response.content))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            if add_border:
                img = ImageOps.expand(img, border=15, fill="black")
            out_stream = io.BytesIO()
            img.save(out_stream, format="PNG")
            out_stream.seek(0)
            return out_stream
    except Exception as e:
        print(f"   [IMG] Exception: {e}")
        return None
    return None


def sanitize_filename(filename: str) -> str:
    return Path(filename).name


def cleanup_job_files(*files):
    for f in files:
        try:
            if os.path.exists(f):
                os.remove(f)
        except Exception as e:
            print(f"Cleanup error for {f}: {e}")


# ==========================================
# BOUNCER HELPERS
# ==========================================
def normalize_bouncer_keywords(keywords_found):
    if isinstance(keywords_found, list):
        return ", ".join(str(keyword).strip() for keyword in keywords_found if str(keyword).strip())
    return str(keywords_found or "").strip()


def build_bouncer_text(title, summary, keywords_found=None):
    return (
        f"Title: {str(title or '').strip()}\n"
        f"Keywords: {normalize_bouncer_keywords(keywords_found)}\n"
        f"Summary: {str(summary or '').strip()}"
    )


def get_bouncer_not_interested_score(
    title,
    summary,
    keywords_found=None,
    profile=DEFAULT_PROFILE,
):
    model_for_profile = get_bouncer_model_for_profile(profile)
    if bouncer_embedder is None or model_for_profile is None:
        return None
    try:
        check_text = build_bouncer_text(title, summary, keywords_found)
        vector = bouncer_embedder.encode([check_text])
        if hasattr(model_for_profile, "predict_proba"):
            probabilities = model_for_profile.predict_proba(vector)[0]
            classes = list(getattr(model_for_profile, "classes_", []))
            for candidate in [0, "0", "not_interested", "not_intrested", "irrelevant", "drop", "dislike"]:
                if candidate in classes:
                    return float(probabilities[classes.index(candidate)])
            print(f"[BOUNCER:{profile}] Could not identify not_interested class: {classes}")
            return None
        prediction = model_for_profile.predict(vector)[0]
        return 1.0 if prediction in [0, "0", "not_interested", "not_intrested", "irrelevant", "drop", "dislike"] else 0.0
    except Exception as e:
        print(f"[BOUNCER:{profile}] Bouncer error: {e}")
        return None


def bouncer_decision(title, summary, keywords_found=None, profile=DEFAULT_PROFILE):
    score = get_bouncer_not_interested_score(title, summary, keywords_found, profile)
    if score is None:
        return {"keep": True, "decision": "keep", "score": None, "reason": f"bouncer_unavailable_{profile}"}
    if score >= BOUNCER_HARD_DROP_THRESHOLD:
        return {"keep": False, "decision": "drop", "score": round(score, 4), "reason": f"high_confidence_not_interested_{profile}"}
    if score >= BOUNCER_LOW_PRIORITY_THRESHOLD:
        return {"keep": True, "decision": "low_priority", "score": round(score, 4), "reason": f"medium_confidence_not_interested_{profile}"}
    return {"keep": True, "decision": "keep", "score": round(score, 4), "reason": f"likely_interesting_{profile}"}


def bouncer_check(title, summary, keywords_found=None, profile=DEFAULT_PROFILE):
    return bouncer_decision(title, summary, keywords_found, profile)["keep"]


# ==========================================
# --- DATA MODELS ---
# ==========================================
class Source(BaseModel):
    name: str


class NewsItem(BaseModel):
    title: str
    master_summary: str
    ppt_summary: Optional[str] = ""
    snippet: Optional[str] = ""
    date: str
    link: str
    top_image: Optional[str] = None
    sources: List[Source] = []
    importance_score: int
    keywords_found: List[str] = []
    region: Optional[str] = "Global"
    full_contents: Optional[str] = ""
    selected_by: Optional[str] = None
    category: Optional[str] = "Tech News"


class ExportRequest(BaseModel):
    items: List[NewsItem]
    filename: str = "SENSE_Brief.pptx"


class VotePayload(BaseModel):
    keywords: list
    summary: str
    vote: str
    title: Optional[str] = ""


# ==========================================
# --- NOT INTERESTED STORE ---
# ==========================================
def load_not_interested_store(request: Request = None):
    profile = get_active_profile_name(request)
    not_interested_file = get_not_interested_file_for_profile(profile)
    if not os.path.exists(not_interested_file):
        return []
    try:
        with open(not_interested_file, "r", encoding="utf-8") as f:
            store = json.load(f)
    except (json.JSONDecodeError, Exception):
        return []

    now = datetime.datetime.now()
    active = []
    expired_count = 0
    for item in store:
        try:
            rejected_at = datetime.datetime.strptime(
                item.get("rejected_at", ""), "%Y-%m-%d %H:%M:%S"
            )
            age_hours = (now - rejected_at).total_seconds() / 3600
            if age_hours <= NOT_INTERESTED_EXPIRY_HOURS:
                active.append(item)
            else:
                expired_count += 1
        except (ValueError, TypeError):
            active.append(item)

    if expired_count > 0:
        try:
            with open(not_interested_file, "w", encoding="utf-8") as f:
                json.dump(active, f, indent=4, ensure_ascii=False)
            print(f"[NOT-INTERESTED:{profile}] Expired {expired_count} entries (>{NOT_INTERESTED_EXPIRY_HOURS}h old)")
        except Exception:
            pass

    return active


def save_not_interested_store(store, request: Request = None):
    profile = get_active_profile_name(request)
    with open(get_not_interested_file_for_profile(profile), "w", encoding="utf-8") as f:
        json.dump(store, f, indent=4, ensure_ascii=False)


def is_already_rejected(title, store):
    normalized_title = title.strip().lower()
    for item in store:
        if item.get("title", "").strip().lower() == normalized_title:
            return True
    return False


# ==========================================
# --- WORKFLOW HELPERS ---
# ==========================================
def load_workflow_store(request: Request = None):
    workflow_file = get_workflow_file_for_request(request)
    if not os.path.exists(workflow_file):
        return {"selected": [], "approved": []}
    try:
        with open(workflow_file, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict):
            return {"selected": [], "approved": []}
        data.setdefault("selected", [])
        data.setdefault("approved", [])
        return data
    except Exception:
        return {"selected": [], "approved": []}


def save_workflow_store(data, request: Request = None):
    with open(get_workflow_file_for_request(request), "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)


# ==========================================
# --- DROPPED ARTICLES DATA ---
# ==========================================
def log_dropped_article(title, summary, keywords_found, bouncer_info=None, profile=DEFAULT_PROFILE):
    dropped_file = os.path.join(ROOT_DIR, "dropped_articles.json")
    bouncer_info = bouncer_info or {}
    new_entry = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "profile": profile,
        "title": title,
        "keyword": keywords_found if keywords_found else [title[:50]],
        "summary": summary,
        "label": "not_interested",
        "bouncer_decision": bouncer_info.get("decision", "drop"),
        "bouncer_score": bouncer_info.get("score"),
        "bouncer_reason": bouncer_info.get("reason", ""),
    }
    with dropped_lock:
        dropped = []
        if os.path.exists(dropped_file):
            try:
                with open(dropped_file, "r", encoding="utf-8") as f:
                    dropped = json.load(f)
            except Exception:
                dropped = []
        dropped.append(new_entry)
        dropped = dropped[-500:]
        with open(dropped_file, "w", encoding="utf-8") as f:
            json.dump(dropped, f, indent=4, ensure_ascii=False)


# ==========================================
# TRAINING AND FILTER HELPERS
# ==========================================
def get_bouncer_summary_from_item(item):
    parts = []
    for key in ["master_summary", "summary", "snippet", "full_content", "full_contents"]:
        text = str(item.get(key, "") or "").strip()
        if text and text not in parts:
            parts.append(text)
    return (" ".join(parts).strip() or str(item.get("title", "")).strip())[:2500]


def run_bouncer_filter_on_items(items, profile=DEFAULT_PROFILE, stage="raw"):
    filtered_items = []
    dropped_count = 0
    low_priority_count = 0
    for item in items if isinstance(items, list) else []:
        decision = bouncer_decision(
            item.get("title", ""),
            get_bouncer_summary_from_item(item),
            item.get("keywords_found", []),
            profile,
        )
        if decision["keep"]:
            item["profile"] = profile
            item["bouncer_stage"] = stage
            item["bouncer_decision"] = decision["decision"]
            item["bouncer_score"] = decision["score"]
            item["bouncer_reason"] = decision["reason"]
            filtered_items.append(item)
            if decision["decision"] == "low_priority":
                low_priority_count += 1
        else:
            dropped_count += 1
            log_dropped_article(
                item.get("title", ""),
                get_bouncer_summary_from_item(item),
                item.get("keywords_found", []),
                decision,
                profile,
            )
    return filtered_items, dropped_count, low_priority_count


def save_training_vote(keywords, summary, vote, title="", profile=DEFAULT_PROFILE):
    training_file = get_training_file_for_profile(profile)
    new_row = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "title": str(title or "").strip(),
        "keyword": keywords,
        "summary": str(summary or "").strip(),
        "label": vote,
        "profile": profile,
    }
    with file_lock:
        memory = []
        if os.path.exists(training_file):
            try:
                with open(training_file, "r", encoding="utf-8") as f:
                    memory = json.load(f)
            except json.JSONDecodeError:
                memory = []

        normalized_title = str(title or "").strip().lower()[:150]
        normalized_new = str(summary or "").strip().lower()[:200]
        normalized_vote = str(vote or "").strip().lower()
        is_duplicate = False
        for existing in memory:
            existing_title = existing.get("title", "").strip().lower()[:150]
            existing_summary = existing.get("summary", "").strip().lower()[:200]
            existing_vote = existing.get("label", "").strip().lower()
            if (
                existing_title == normalized_title
                and existing_summary == normalized_new
                and existing_vote == normalized_vote
            ):
                is_duplicate = True
                break

        if not is_duplicate:
            memory.append(new_row)
            with open(training_file, "w", encoding="utf-8") as f:
                json.dump(memory, f, indent=4, ensure_ascii=False)
        else:
            print(f"Dedup: Skipped duplicate {vote} vote for: {summary[:50]}...")

    return len(memory)


def reload_bouncer_model_for_profile(profile=DEFAULT_PROFILE):
    global bouncer_model
    model_file = get_bouncer_model_file_for_profile(profile)
    if not os.path.exists(model_file):
        return False
    try:
        with open(model_file, "rb") as f:
            bouncer_models[profile] = pickle.load(f)
        if profile == DEFAULT_PROFILE:
            bouncer_model = bouncer_models.get(profile)
        return True
    except Exception as e:
        print(f"[BOUNCER:{profile}] Failed to reload model: {e}")
        return False


def retrain_and_reload(profile=DEFAULT_PROFILE):
    if not train_lock.acquire(blocking=False):
        print(f"[BOUNCER:{profile}] AI is already training. Skipping duplicate request.")
        return
    try:
        print(f"\n[BOUNCER:{profile}] Retraining with new data...")
        subprocess.run(
            [sys.executable, "train_bouncer.py", "--profile", profile],
            check=True,
        )
        reload_bouncer_model_for_profile(profile)
        print(f"[BOUNCER:{profile}] Brain successfully upgraded and reloaded.\n")
    except Exception as e:
        print(f"[BOUNCER:{profile}] Failed to retrain Bouncer: {e}")
    finally:
        train_lock.release()


# ==========================================
# --- USAGE TRACKING HELPERS ---
# ==========================================
def load_tracker():
    if not os.path.exists(USAGE_TRACKER_FILE):
        return {}
    try:
        with open(USAGE_TRACKER_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except:
        return {}


def save_tracker(data):
    with open(USAGE_TRACKER_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)


def get_device_id(ip, fingerprint):
    raw = f"{ip}_{fingerprint}"
    return hashlib.md5(raw.encode()).hexdigest()[:12]


def get_today():
    return datetime.datetime.now().strftime("%Y-%m-%d")


def get_empty_day():
    return {
        "page_loads": 0,
        "searches": [],
        "articles_clicked": 0,
        "votes_interested": 0,
        "votes_not_interested": 0,
        "exports": [],
        "briefing_views": 0,
        "heartbeats": 0,
        "voc_feedback": [],
    }


def purge_old_entries(device_data, keep_days=30):
    cutoff = (
        datetime.datetime.now() - datetime.timedelta(days=keep_days)
    ).strftime("%Y-%m-%d")
    activity = device_data.get("activity", {})
    device_data["activity"] = {
        date: data for date, data in activity.items()
        if date >= cutoff
    }
    return device_data


# ==========================================
# Legacy default-only scheduler retained for archive compatibility; the
# profile-aware scheduler definition below is the runtime entry point.
# ==========================================
def _legacy_run_morning_briefing():
    global SCHEDULER_STATUS

    with scheduler_lock:
        if SCHEDULER_STATUS["is_active"]:
            print("[SCHEDULER] Already running. Skipping.")
            return

        running_manual = sum(
            1 for j in active_jobs.values() if j.get("status") == "running"
        )
        if running_manual > 0:
            print(
                f"[SCHEDULER] {running_manual} manual scan(s) in progress. "
                f"Deferring autonomous run by 10 minutes."
            )
            threading.Timer(600, _legacy_run_morning_briefing).start()
            return

        SCHEDULER_STATUS["is_active"] = True
        SCHEDULER_STATUS["message"] = "Autonomous Engine Scanning..."
        SCHEDULER_STATUS["mode"] = "autonomous"

    scheduler_job_id = (
        f"scheduler_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    )

    try:
        today = datetime.datetime.now()
        yesterday = today - datetime.timedelta(days=1)
        from_date = yesterday.strftime("%Y-%m-%d")
        to_date = today.strftime("%Y-%m-%d")

        root_dir = os.getcwd()
        output_file = os.path.join(
            root_dir, f"ui_results_{scheduler_job_id}.json"
        )
        cluster_file = os.path.join(
            root_dir, f"clustered_results_{scheduler_job_id}.json"
        )

        for pattern in [
            "ui_results_scheduler_*.json",
            "clustered_results_scheduler_*.json",
        ]:
            for old_file in glob.glob(os.path.join(root_dir, pattern)):
                try:
                    os.remove(old_file)
                except:
                    pass

        print(f"[SCHEDULER] Starting scan: {scheduler_job_id}")
        print(
            f"   -> Deploying Spiders for keywords: {MORNING_KEYWORDS[:50]}..."
        )

        cmd = [
            sys.executable, "-m", "scrapy", "crawl", "news_spider",
            "-a", f"keyword={MORNING_KEYWORDS}",
            "-a", f"from_date={from_date}",
            "-a", f"to_date={to_date}",
            "-a", "target_sites=All",
            "-s", "ROBOTSTXT_OBEY=False",
            "-s", "USER_AGENT=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "-s", "DNS_RESOLVER=scrapy.resolver.CachingThreadedResolver",
            "-s", "TWISTED_REACTOR=twisted.internet.asyncioreactor.AsyncioSelectorReactor",
            "-O", output_file,
        ]
        spider_cwd = os.path.join(root_dir, "news_aggregator")
        subprocess.run(cmd, cwd=spider_cwd, timeout=3600)

        if os.path.exists(output_file):
            try:
                print("   -> [SCHEDULER] Running AI Gatekeeper...")
                with open(output_file, "r", encoding="utf-8") as f:
                    raw_data = json.load(f)

                filtered_data = []
                dropped_count = 0

                for item in raw_data:
                    title = item.get("title", "")
                    summary = item.get("snippet", item.get("summary", ""))
                    keywords_found = item.get("keywords_found", [])

                    if bouncer_check(title, summary, keywords_found):
                        filtered_data.append(item)
                    else:
                        dropped_count += 1
                        print(f"   Dropped: {title}", flush=True)
                        log_dropped_article(title, summary, keywords_found)

                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(filtered_data, f, indent=4)

                print(
                    f"   -> [SCHEDULER] Gatekeeper removed {dropped_count} articles.",
                    flush=True,
                )
            except Exception as e:
                print(
                    f"   -> [SCHEDULER] Bouncer error, skipping filter: {e}"
                )

        print("   -> Activating Semantic Fusion Engine (Fast Mode)...")
        subprocess.run(
            [
                sys.executable, "semantic_clustering.py",
                "--job-id", scheduler_job_id,
                "--fast-mode",
            ],
            cwd=root_dir,
            timeout=3600,
        )

        if os.path.exists(cluster_file):
            with open(cluster_file, "r", encoding="utf-8") as f:
                results = json.load(f)

            if results:
                for r in results:
                    r["category"] = assign_category(
                        r.get("title", ""),
                        r.get("master_summary", "") or r.get("snippet", ""),
                    )
                    r.update(apply_learned_region(r))

                print(f"   -> Archiving {len(results)} intelligence items...")
                learner.log_search_data(MORNING_KEYWORDS, results)

                timestamp = today.strftime("%Y-%m-%d_%H-%M-%S")
                history_path = os.path.join(
                    HISTORY_DIR, f"briefing_{timestamp}.json"
                )
                with open(history_path, "w", encoding="utf-8") as f:
                    json.dump(results, f, indent=4)

                try:
                    df = pd.DataFrame(results)
                    df["search_timestamp"] = today
                    df["search_keywords"] = "MORNING_BRIEFING"
                    df["global_week"] = today.isocalendar()[1]

                    if os.path.exists(MANUAL_LOG_FILE):
                        with pd.ExcelWriter(
                            MANUAL_LOG_FILE,
                            mode="a",
                            if_sheet_exists="overlay",
                            engine="openpyxl",
                        ) as writer:
                            pd.concat(
                                [pd.read_excel(MANUAL_LOG_FILE), df],
                                ignore_index=True,
                            ).to_excel(writer, index=False)
                    else:
                        df.to_excel(MANUAL_LOG_FILE, index=False)
                except Exception as e:
                    print(f"Excel Logging Error: {e}")

                print("[SCHEDULER] Morning Briefing Complete & Ready.")
            else:
                print("[SCHEDULER] Briefing finished but no news was found.")
        else:
            print("[SCHEDULER] Failed to generate cluster file.")

    except subprocess.TimeoutExpired:
        print("\n[SCHEDULER ERROR] Process timed out and was terminated.")
    except Exception as e:
        print(f"\n[SCHEDULER ERROR] Critical failure: {e}")
    finally:
        with scheduler_lock:
            SCHEDULER_STATUS["is_active"] = False
            SCHEDULER_STATUS["message"] = "Morning Briefing Complete."
            SCHEDULER_STATUS["mode"] = "idle"
        cleanup_job_files(output_file, cluster_file)
        print(f"[SCHEDULER] Cleaned up temp files for {scheduler_job_id}")


# ==========================================
# PROFILE-AWARE AUTONOMOUS SCHEDULER
# ==========================================
def run_scheduler_for_profile(profile: str):
    config = get_profile_config(profile)
    today = datetime.datetime.now()
    scheduler_job_id = f"scheduler_{profile}_{today.strftime('%Y%m%d_%H%M%S')}"
    output_file = os.path.join(ROOT_DIR, f"ui_results_{scheduler_job_id}.json")
    cluster_file = os.path.join(ROOT_DIR, f"clustered_results_{scheduler_job_id}.json")
    history_dir = get_profile_history_dir(profile)

    if not os.path.exists(config["sites_file"]):
        print(f"[SCHEDULER:{profile}] Missing sites file: {config['sites_file']}")
        return

    try:
        from_date = (today - datetime.timedelta(days=1)).strftime("%Y-%m-%d")
        to_date = today.strftime("%Y-%m-%d")
        cmd = [
            sys.executable, "-m", "scrapy", "crawl", "news_spider",
            "-a", f"keyword={config['keywords']}",
            "-a", f"from_date={from_date}",
            "-a", f"to_date={to_date}",
            "-a", "target_sites=All",
            "-a", f"sites_file={config['sites_file']}",
            "-s", "ROBOTSTXT_OBEY=False",
            "-s", "USER_AGENT=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "-s", "DNS_RESOLVER=scrapy.resolver.CachingThreadedResolver",
            "-s", "TWISTED_REACTOR=twisted.internet.asyncioreactor.AsyncioSelectorReactor",
            "-O", output_file,
        ]
        print(f"[SCHEDULER:{profile}] Starting scan: {scheduler_job_id}")
        subprocess.run(cmd, cwd=os.path.join(ROOT_DIR, "news_aggregator"), timeout=3600)

        if os.path.exists(output_file) and config["use_bouncer"]:
            with open(output_file, "r", encoding="utf-8") as f:
                raw_data = json.load(f)
            raw_data, dropped, low_priority = run_bouncer_filter_on_items(
                raw_data, profile, "scheduler_raw"
            )
            with open(output_file, "w", encoding="utf-8") as f:
                json.dump(raw_data, f, indent=4, ensure_ascii=False)
            print(f"[SCHEDULER:{profile}] Dropped {dropped}; low priority kept {low_priority}.")

        subprocess.run(
            [sys.executable, "semantic_clustering.py", "--job-id", scheduler_job_id, "--fast-mode"],
            cwd=ROOT_DIR,
            timeout=3600,
        )
        if not os.path.exists(cluster_file):
            print(f"[SCHEDULER:{profile}] Failed to generate cluster file.")
            return
        with open(cluster_file, "r", encoding="utf-8") as f:
            results = json.load(f)
        if results and config["use_bouncer"]:
            results, dropped, low_priority = run_bouncer_filter_on_items(
                results, profile, "scheduler_final"
            )
        for item in results:
            item["profile"] = profile
            item["category"] = assign_category(
                item.get("title", ""),
                item.get("master_summary", "") or item.get("snippet", ""),
            )
            item.update(apply_learned_region(item, profile))
        if not results:
            print(f"[SCHEDULER:{profile}] Finished but no news was found.")
            return
        learner.log_search_data(config["keywords"], results)
        timestamp = today.strftime("%Y-%m-%d_%H-%M-%S")
        history_path = os.path.join(history_dir, f"briefing_{timestamp}.json")
        with open(history_path, "w", encoding="utf-8") as f:
            json.dump(results, f, indent=4, ensure_ascii=False)
        print(f"[SCHEDULER:{profile}] Briefing complete: {history_path}")
    except subprocess.TimeoutExpired:
        print(f"[SCHEDULER:{profile}] Process timed out and was terminated.")
    except Exception as e:
        print(f"[SCHEDULER:{profile}] Critical failure: {e}")
    finally:
        cleanup_job_files(output_file, cluster_file)


def run_morning_briefing():
    global SCHEDULER_STATUS
    with scheduler_lock:
        if SCHEDULER_STATUS["is_active"]:
            return
        if any(
            job.get("status") in {"queued", "running"}
            for job in active_jobs.values()
        ):
            print("[SCHEDULER] Manual scan admitted first. Deferring autonomous scan.")
            threading.Timer(600, run_morning_briefing).start()
            return
        SCHEDULER_STATUS.update(
            {"is_active": True, "message": "Autonomous Engine Scanning...", "mode": "autonomous"}
        )
    try:
        ensure_profile_storage()
        for profile in [DEFAULT_PROFILE, BROADCAST_PROFILE]:
            SCHEDULER_STATUS["message"] = f"Autonomous Engine Scanning {profile} profile..."
            run_scheduler_for_profile(profile)
    finally:
        with scheduler_lock:
            SCHEDULER_STATUS.update(
                {"is_active": False, "message": "Morning Briefing Complete.", "mode": "idle"}
            )


# ==========================================
# --- LIFECYCLE ---
# ==========================================
@asynccontextmanager
async def lifespan(app: FastAPI):
    ensure_profile_storage()
    scheduler = BackgroundScheduler()
    next_run = datetime.datetime.now() + datetime.timedelta(minutes=2)

    list_of_files = [
        latest
        for profile in [DEFAULT_PROFILE, BROADCAST_PROFILE]
        if (latest := get_latest_briefing_file_for_profile(profile))
    ]
    if list_of_files:
        latest_file = max(list_of_files, key=os.path.getmtime)
        try:
            last_run_time = datetime.datetime.fromtimestamp(os.path.getmtime(latest_file))
            target_next_run = last_run_time + datetime.timedelta(hours=4)
            if target_next_run > datetime.datetime.now():
                next_run = target_next_run
                print(f"[SCHEDULER] Fresh intelligence. Next scan: {next_run.strftime('%I:%M %p')}.")
            else:
                print(f"[SCHEDULER] Intelligence stale. Scan at {next_run.strftime('%I:%M %p')} (2 min startup delay).")
        except Exception as e:
            print(f"[SCHEDULER] Error reading history ({e}). Scan in 2 minutes.")
    else:
        print(f"[SCHEDULER] No history found. Initial scan at {next_run.strftime('%I:%M %p')}.")

    scheduler.add_job(run_morning_briefing, "interval", hours=4, next_run_time=next_run)
    scheduler.start()
    print("SYSTEM: Autonomous Intelligence Engine online.")
    yield
    scheduler.shutdown()
    ml_executor.shutdown(wait=False)


# ==========================================
# --- APP INIT ---
# ==========================================
app = FastAPI(lifespan=lifespan)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ==========================================
# --- STATIC FILES ---
# ==========================================
abs_frontend_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "..", "frontend", "dist"
)
if os.path.exists(os.path.join(abs_frontend_path, "assets")):
    app.mount(
        "/assets",
        StaticFiles(directory=os.path.join(abs_frontend_path, "assets")),
        name="assets",
    )


# ==========================================
# PROFILE ENDPOINT
# ==========================================
@app.get("/profile")
def get_profile(request: Request):
    profile = get_active_profile_name(request)
    config = get_profile_config(profile)
    return {
        "status": "success",
        "ip": get_client_ip(request),
        "profile": profile,
        "label": config["label"],
        "is_broadcast": profile == BROADCAST_PROFILE,
        "paths": {
            "sites_file": config["sites_file"],
            "history_dir": config["history_dir"],
            "workflow_file": get_workflow_file_for_request(request),
            "not_interested_file": get_not_interested_file_for_profile(profile),
            "training_file": get_training_file_for_profile(profile),
            "bouncer_model_file": get_bouncer_model_file_for_profile(profile),
        },
        "all_profiles": {
            DEFAULT_PROFILE: get_profile_debug_info(DEFAULT_PROFILE),
            BROADCAST_PROFILE: get_profile_debug_info(BROADCAST_PROFILE),
        },
    }


# ==========================================
# --- EXPORT: POWERPOINT ---
# ==========================================
@app.post("/export-ppt")
async def export_ppt(http_request: Request, request: ExportRequest):
    if get_active_profile_name(http_request) == BROADCAST_PROFILE:
        raise HTTPException(
            status_code=403,
            detail="PowerPoint export is disabled for Broadcast profile.",
        )
    safe_filename = sanitize_filename(request.filename)
    TEMPLATE_PATH = "template.pptx"
    if not os.path.exists(TEMPLATE_PATH):
        raise HTTPException(status_code=404, detail="template.pptx not found")

    prs = Presentation(TEMPLATE_PATH)
    cover_layout = None
    news_layout = None
    for layout in prs.slide_layouts:
        if layout.name == "CoverLayout":
            cover_layout = layout
        if layout.name == "NewsLayout":
            news_layout = layout

    if not cover_layout:
        cover_layout = prs.slide_layouts[0]
    if not news_layout:
        news_layout = prs.slide_layouts[0] if len(prs.slide_layouts) == 1 else prs.slide_layouts[1]

    ph_map = {"title": -1, "summary": -1, "link": -1, "insight": -1, "picture": -1, "date": -1, "team": -1}
    for shape in news_layout.placeholders:
        idx = shape.placeholder_format.idx
        if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            ph_map["picture"] = idx
            continue
        if shape.has_text_frame:
            text = shape.text.strip()
            if "#TITLE" in text: ph_map["title"] = idx
            elif "#SUMMARY" in text: ph_map["summary"] = idx
            elif "#LINK" in text: ph_map["link"] = idx
            elif "#INSIGHT" in text: ph_map["insight"] = idx
            elif "#DATE_HERE" in text: ph_map["date"] = idx
            elif "#Targated_SRID_Team" in text: ph_map["team"] = idx

    slide = prs.slides.add_slide(cover_layout)
    for shape in slide.shapes:
        if shape.has_text_frame and "#DATE_HERE" in shape.text:
            shape.text = datetime.datetime.now().strftime("%b'%y")

    for item in request.items:
        slide = prs.slides.add_slide(news_layout)
        ai_opinion = generate_opinion(item.master_summary)
        target_team = determine_target_team(item.title, item.master_summary)

        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == ph_map["title"]:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                kw = item.keywords_found[0] if item.keywords_found else ""
                if kw:
                    parts = re.split(f"({re.escape(kw)})", item.title, flags=re.IGNORECASE)
                    for part in parts:
                        run = p.add_run()
                        run.text = part
                        if part.lower() == kw.lower():
                            run.font.bold = True
                            run.font.underline = True
                            run.font.color.rgb = RGBColor(0, 112, 192)
                        else:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0, 0, 0)
                else:
                    run = p.add_run()
                    run.text = item.title
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)
            elif idx == ph_map["summary"]:
                tf = shape.text_frame
                tf.clear()
                target_summary = item.ppt_summary if item.ppt_summary else item.master_summary
                sentences = [s.strip() + "." for s in target_summary.split(". ") if s.strip()]
                if not sentences: sentences = [target_summary]
                for i, sentence in enumerate(sentences):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = f"• {sentence}"
                    p.font.name = "Calibri"
                    p.font.size = Pt(18)
            elif idx == ph_map["link"]:
                shape.text = item.link
                try: shape.text_frame.paragraphs[0].font.size = Pt(10)
                except: pass
            elif idx == ph_map["insight"]:
                shape.text = f"Insight : {ai_opinion}"
                try: shape.text_frame.paragraphs[0].font.size = Pt(14)
                except: pass
            elif idx == ph_map["date"]:
                shape.text = datetime.datetime.now().strftime("%b'%y")
            elif idx == ph_map["team"]:
                shape.text = f"Targeted SRID TEAM : {target_team}"
            elif idx == ph_map["picture"] or shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                img_stream = download_image_for_export(item.top_image, add_border=False)
                if img_stream:
                    try:
                        left, top = shape.left, shape.top
                        width, height = shape.width, shape.height
                        pic = slide.shapes.add_picture(img_stream, left, top, width, height)
                        try: pic.click_action.hyperlink.address = item.link
                        except: pass
                    except Exception as e:
                        print(f"   [PPT] Insert Error: {e}")

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={safe_filename}"},
    )


# ==========================================
# --- EXPORT: EXCEL ---
# ==========================================
@app.post("/export-excel")
async def export_excel(request: ExportRequest):
    safe_filename = sanitize_filename(request.filename)
    wb = Workbook()
    ws = wb.active
    ws.title = "SENSE Report"
    ws["A2"] = "Weekly Report"
    ws["A2"].font = Font(size=20, bold=True)

    headers = ["Sr. No.", "WK", "Date", "Name of the Initiator", "Category", "Keyword", "News Highlight", "URL"]
    header_fill = PatternFill(start_color="ADD8E6", end_color="ADD8E6", fill_type="solid")
    for col_num, header_title in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col_num, value=header_title)
        cell.font = Font(bold=True)
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")

    for idx, item in enumerate(request.items, 1):
        row_num = 4 + idx
        try:
            dt = datetime.datetime.strptime(item.date, "%Y-%m-%d")
            wk = dt.isocalendar()[1]
        except:
            wk = datetime.datetime.now().isocalendar()[1]

        initiator = getattr(item, "selected_by", "")
        category_val = getattr(item, "category", "Tech News")
        kw = ", ".join(item.keywords_found) if item.keywords_found else ""
        summary_sentences = [s.strip() + "." for s in item.master_summary.split(". ") if s.strip()]
        bullets = "\n".join([f"• {s}" for s in summary_sentences])
        highlight_text = f"{item.title.upper()}\n\n{bullets}"

        ws.cell(row=row_num, column=1, value=idx).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=2, value=wk).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=3, value=item.date).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=4, value=initiator).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=5, value=category_val.upper()).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=6, value=kw).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=7, value=highlight_text).alignment = Alignment(wrap_text=True, vertical="top")
        ws.cell(row=row_num, column=8, value=item.link).alignment = Alignment(vertical="top")

    ws.column_dimensions["A"].width = 8
    ws.column_dimensions["B"].width = 6
    ws.column_dimensions["C"].width = 12
    ws.column_dimensions["D"].width = 22
    ws.column_dimensions["E"].width = 18
    ws.column_dimensions["F"].width = 20
    ws.column_dimensions["G"].width = 80
    ws.column_dimensions["H"].width = 40

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={safe_filename}"},
    )


# ==========================================
# --- EXPORT: WORD ---
# ==========================================
@app.post("/export-word")
async def export_word(request: ExportRequest):
    safe_filename = sanitize_filename(request.filename)
    doc = Document()
    section = doc.sections[0]
    section.page_width = DocxInches(8.5)
    section.page_height = DocxInches(11)
    section.left_margin = DocxInches(1)
    section.right_margin = DocxInches(1)
    section.top_margin = DocxInches(1)
    section.bottom_margin = DocxInches(1)
    section.different_first_page_header_footer = True

    p_brand = doc.add_paragraph()
    run_brand = p_brand.add_run("Samsung General")
    run_brand.font.name = "Calibri"
    run_brand.font.size = DocxPt(7)
    run_brand.font.color.rgb = DocxRGBColor(0, 160, 70)
    p_brand.paragraph_format.space_before = DocxPt(0)
    p_brand.paragraph_format.space_after = DocxPt(0)

    header = section.header
    header_para = header.paragraphs[0]
    header_para.clear()
    run_h_left = header_para.add_run("SENSE Intelligence Brief")
    run_h_left.font.name = "Calibri"
    run_h_left.font.size = DocxPt(9)
    run_h_left.font.color.rgb = DocxRGBColor(120, 120, 120)
    header_para.add_run("\t")
    run_h_right = header_para.add_run(datetime.datetime.now().strftime("%d %B %Y"))
    run_h_right.font.name = "Calibri"
    run_h_right.font.size = DocxPt(9)
    run_h_right.font.color.rgb = DocxRGBColor(120, 120, 120)

    pPr_h = header_para._p.get_or_add_pPr()
    tabs_h = OxmlElement("w:tabs")
    tab_h = OxmlElement("w:tab")
    tab_h.set(qn("w:val"), "right")
    tab_h.set(qn("w:pos"), "9360")
    tabs_h.append(tab_h)
    pPr_h.append(tabs_h)
    pBdr_h = OxmlElement("w:pBdr")
    bottom_h = OxmlElement("w:bottom")
    bottom_h.set(qn("w:val"), "single")
    bottom_h.set(qn("w:sz"), "4")
    bottom_h.set(qn("w:space"), "1")
    bottom_h.set(qn("w:color"), "CCCCCC")
    pBdr_h.append(bottom_h)
    pPr_h.append(pBdr_h)

    footer = section.footer
    footer_para = footer.paragraphs[0]
    footer_para.clear()
    run_conf = footer_para.add_run("CONFIDENTIAL - For Internal Use Only")
    run_conf.font.name = "Calibri"
    run_conf.font.size = DocxPt(8)
    run_conf.font.color.rgb = DocxRGBColor(150, 150, 150)
    footer_para.add_run("\t")
    run_pg = footer_para.add_run()
    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")
    instrText = OxmlElement("w:instrText")
    instrText.text = "PAGE"
    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")
    run_pg._r.append(fldChar1)
    run_pg._r.append(instrText)
    run_pg._r.append(fldChar2)
    run_pg.font.name = "Calibri"
    run_pg.font.size = DocxPt(8)
    run_pg.font.color.rgb = DocxRGBColor(150, 150, 150)

    pPr_f = footer_para._p.get_or_add_pPr()
    tabs_f = OxmlElement("w:tabs")
    tab_f = OxmlElement("w:tab")
    tab_f.set(qn("w:val"), "right")
    tab_f.set(qn("w:pos"), "9360")
    tabs_f.append(tab_f)
    pPr_f.append(tabs_f)
    pBdr_f = OxmlElement("w:pBdr")
    top_f = OxmlElement("w:top")
    top_f.set(qn("w:val"), "single")
    top_f.set(qn("w:sz"), "4")
    top_f.set(qn("w:space"), "1")
    top_f.set(qn("w:color"), "CCCCCC")
    pBdr_f.append(top_f)
    pPr_f.append(pBdr_f)

    def add_section_label(doc, text):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.name = "Calibri"
        run.font.size = DocxPt(11)
        run.font.bold = True
        run.font.color.rgb = DocxRGBColor(0, 51, 102)
        run.font.all_caps = True
        p.paragraph_format.space_before = DocxPt(12)
        p.paragraph_format.space_after = DocxPt(4)
        return p

    def add_rule(doc):
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "4")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "CCCCCC")
        pBdr.append(bottom)
        pPr.append(pBdr)
        p.paragraph_format.space_before = DocxPt(0)
        p.paragraph_format.space_after = DocxPt(10)

    for item in request.items:
        doc.add_page_break()
        p_title = doc.add_paragraph()
        p_title.paragraph_format.space_after = DocxPt(4)
        kw = item.keywords_found[0] if item.keywords_found else ""
        if kw:
            parts = re.split(f"({re.escape(kw)})", item.title, flags=re.IGNORECASE)
            for part in parts:
                run = p_title.add_run(part)
                run.font.name = "Calibri"
                run.font.size = DocxPt(16)
                run.font.bold = True
                run.font.color.rgb = DocxRGBColor(0, 0, 0)
        else:
            run = p_title.add_run(item.title)
            run.font.name = "Calibri"
            run.font.size = DocxPt(16)
            run.font.bold = True
            run.font.color.rgb = DocxRGBColor(0, 0, 0)

        region_text = (getattr(item, "region", None) or "Global").upper()
        importance = getattr(item, "importance_score", 0)
        p_meta = doc.add_paragraph()
        run_meta = p_meta.add_run(f"{item.date}     |     Region: {region_text}     |     Importance: {importance}/10")
        run_meta.font.name = "Calibri"
        run_meta.font.size = DocxPt(10)
        run_meta.font.color.rgb = DocxRGBColor(130, 130, 130)
        p_meta.paragraph_format.space_after = DocxPt(6)

        add_rule(doc)

        img_stream = download_image_for_export(item.top_image, add_border=True)
        if img_stream:
            p_img = doc.add_paragraph()
            p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_img.add_run().add_picture(img_stream, width=DocxInches(5.5))
            p_img.paragraph_format.space_after = DocxPt(10)

        add_section_label(doc, "Intelligence Summary")
        sentences = [s.strip() + "." for s in item.master_summary.split(". ") if s.strip()]
        if not sentences: sentences = [item.master_summary]
        for sentence in sentences:
            p_b = doc.add_paragraph(style="List Bullet")
            run_b = p_b.add_run(sentence)
            run_b.font.name = "Calibri"
            run_b.font.size = DocxPt(11)
            run_b.font.color.rgb = DocxRGBColor(30, 30, 30)

        ai_opinion = generate_opinion(item.master_summary)
        add_section_label(doc, "Analysis")
        p_insight = doc.add_paragraph()
        p_insight.paragraph_format.left_indent = DocxInches(0.3)
        p_insight.paragraph_format.space_after = DocxPt(4)
        run_i = p_insight.add_run(ai_opinion)
        run_i.font.name = "Calibri"
        run_i.font.size = DocxPt(11)
        run_i.font.italic = True
        run_i.font.color.rgb = DocxRGBColor(80, 80, 80)

        team = determine_target_team(item.title, item.master_summary)
        add_section_label(doc, "Routing")
        p_team = doc.add_paragraph()
        run_team = p_team.add_run(team)
        run_team.font.name = "Calibri"
        run_team.font.size = DocxPt(11)
        run_team.font.color.rgb = DocxRGBColor(30, 30, 30)

        add_section_label(doc, "Sources")
        sources = item.sources if item.sources else []
        if sources:
            for source in sources:
                p_src = doc.add_paragraph()
                run_src = p_src.add_run(source.name)
                run_src.font.name = "Calibri"
                run_src.font.size = DocxPt(10)
                run_src.font.color.rgb = DocxRGBColor(60, 60, 60)
                p_src.paragraph_format.space_after = DocxPt(2)
        else:
            p_src = doc.add_paragraph()
            run_src = p_src.add_run(getattr(item, "source", "Unknown"))
            run_src.font.name = "Calibri"
            run_src.font.size = DocxPt(10)
            run_src.font.color.rgb = DocxRGBColor(60, 60, 60)

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={safe_filename}"},
    )


# ==========================================
# --- WORKFLOW ENDPOINTS ---
# ==========================================
@app.get("/workflow")
def get_workflow(request: Request):
    profile = get_profile_for_request(request)
    store = load_workflow_store(request)
    return {
        "selected": apply_learned_regions(store.get("selected", []), profile),
        "approved": apply_learned_regions(store.get("approved", []), profile),
        "profile": profile,
    }


@app.post("/workflow/select")
def select_news(request: Request, item: dict = Body(...)):
    profile = get_profile_for_request(request)
    store = load_workflow_store(request)
    if any(i.get("title") == item.get("title") for i in store["selected"]):
        return {"status": "exists", "message": "Already selected", "profile": profile}
    item["profile"] = profile
    store["selected"].append(item)
    save_workflow_store(store, request)
    return {"status": "success", "count": len(store["selected"]), "profile": profile}


@app.post("/workflow/approve")
def approve_news(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    item_title = payload.get("title")
    key = payload.get("key")
    if key != DIRECTOR_KEY:
        return {"status": "error", "message": "Invalid Director Key", "profile": profile}
    store = load_workflow_store(request)
    item_to_approve = next((i for i in store["selected"] if i["title"] == item_title), None)
    if not item_to_approve:
        return {"status": "error", "message": "Item not found", "profile": profile}
    store["selected"] = [i for i in store["selected"] if i["title"] != item_title]
    item_to_approve["profile"] = profile
    store["approved"].append(item_to_approve)
    save_workflow_store(store, request)
    return {"status": "success", "message": "Approved", "profile": profile}


@app.post("/workflow/remove")
def remove_news(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    title = payload.get("title")
    list_type = payload.get("list_type")
    if list_type not in ["selected", "approved"]:
        return {"status": "error", "message": "Invalid list type", "profile": profile}
    store = load_workflow_store(request)
    store[list_type] = [i for i in store[list_type] if i["title"] != title]
    save_workflow_store(store, request)
    return {"status": "success", "profile": profile}


# ==========================================
# --- REGION CORRECTION / LEARNING ENDPOINT ---
# ==========================================
@app.post("/region/correct")
def correct_region(request: Request, payload: dict = Body(...)):
    profile = get_active_profile_name(request)
    title = str(payload.get("title", "")).strip()
    region = normalize_region_label(payload.get("region"))
    keywords = normalize_region_keywords(payload.get("keywords", []))
    reason = str(payload.get("reason", "")).strip()
    if not title:
        raise HTTPException(status_code=400, detail="Article title is required")
    if not region:
        raise HTTPException(status_code=400, detail="Region must be Local or Global")
    if not keywords:
        raise HTTPException(
            status_code=400,
            detail="Add at least one keyword so future scans can learn this correction",
        )

    other_region = "Global" if region == "Local" else "Local"
    with region_learning_lock:
        learned = load_region_learning(profile)
        learned[other_region] = [
            keyword for keyword in learned[other_region] if keyword not in keywords
        ]
        for keyword in keywords:
            if keyword not in learned[region]:
                learned[region].append(keyword)
        learned["corrections"].append({
            "title": title,
            "previous_region": normalize_region_label(payload.get("previous_region")),
            "region": region,
            "keywords": keywords,
            "reason": reason,
            "created_at": datetime.datetime.now().isoformat(timespec="seconds"),
        })
        learned["corrections"] = learned["corrections"][-500:]
        save_region_learning(learned, profile)

    return {
        "status": "success",
        "profile": profile,
        "region": region,
        "keywords": keywords,
        "message": f"Saved. Future scans will use {len(keywords)} learned keyword(s) for {region} signals.",
    }


# ==========================================
# --- SITES ENDPOINTS ---
# ==========================================
@app.get("/sites")
def get_sites(request: Request):
    profile = get_profile_for_request(request)
    sites_path = get_sites_file_for_profile(profile)
    if os.path.exists(sites_path):
        with open(sites_path, "r", encoding="utf-8") as f:
            sites = json.load(f)
            sites.sort(key=lambda x: x.get("name", "").lower())
            return sites
    return []


@app.post("/sites")
def add_site(site: dict, request: Request):
    profile = get_profile_for_request(request)
    sites_path = get_sites_file_for_profile(profile)
    sites = []
    if os.path.exists(sites_path):
        with open(sites_path, "r", encoding="utf-8") as f:
            sites = json.load(f)
    sites.append(site)
    with open(sites_path, "w", encoding="utf-8") as f:
        json.dump(sites, f, indent=4)
    return {"status": "success", "profile": profile, "sites_file": sites_path}


# ==========================================
# --- STATUS ENDPOINT ---
# ==========================================
@app.get("/status")
def get_system_status(request: Request):
    active_profile = get_active_profile_name(request)
    with scheduler_lock:
        running_jobs = sum(
            1
            for job in active_jobs.values()
            if job.get("status") in {"queued", "running"}
        )
        return {
            **SCHEDULER_STATUS,
            "active_profile": active_profile,
            "active_manual_jobs": running_jobs,
            "capacity_remaining": crawl_semaphore._value,
            "profiles": {
                DEFAULT_PROFILE: get_profile_debug_info(DEFAULT_PROFILE),
                BROADCAST_PROFILE: get_profile_debug_info(BROADCAST_PROFILE),
            },
        }


# ==========================================
# DOSSIER INSIGHT ENDPOINT
# ==========================================
@app.post("/insight")
def get_dossier_insight(request: Request, item: dict = Body(...)):
    profile = get_active_profile_name(request)
    insight, generated_by = generate_why_it_matters(item, profile)
    return {
        "status": "success",
        "profile": profile,
        "why_matters": insight,
        "generated_by": generated_by,
    }


# ==========================================
# --- BRIEFING META ENDPOINT ---
# ==========================================
@app.get("/briefing/meta")
def get_briefing_meta(request: Request):
    profile = get_profile_for_request(request)
    latest = get_latest_briefing_file_for_profile(profile)
    if not latest:
        return {"last_updated": None, "count": 0, "filename": None, "profile": profile}
    try:
        with open(latest, "r", encoding="utf-8") as f:
            data = json.load(f)
        return {
            "last_updated": datetime.datetime.fromtimestamp(os.path.getmtime(latest)).isoformat(),
            "filename": os.path.basename(latest),
            "count": len(data),
            "profile": profile,
        }
    except Exception:
        return {"last_updated": None, "count": 0, "filename": None, "profile": profile}


# ==========================================
# --- LATEST BRIEFING ENDPOINT ---
# ==========================================
@app.get("/latest-briefing")
def get_latest_briefing(request: Request):
    profile = get_profile_for_request(request)
    latest = get_latest_briefing_file_for_profile(profile)
    if latest:
        try:
            with open(latest, "r", encoding="utf-8") as file_obj:
                data = json.load(file_obj)
                if data and len(data) > 0:
                    return {
                        "status": "success",
                        "result": apply_learned_regions(data, profile),
                        "type": "scheduler",
                        "source": "shared",
                        "profile": profile,
                        "filename": os.path.basename(latest),
                        "generated_at": datetime.datetime.fromtimestamp(os.path.getmtime(latest)).strftime("%d %b %Y, %I:%M %p"),
                    }
        except Exception as e:
            return {"status": "error", "result": [], "profile": profile, "message": str(e)}
    return {"status": "empty", "result": [], "profile": profile}


# ==========================================
# --- BRIEFING REMOVE/RESTORE (Disk) ---
# ==========================================
@app.post("/briefing/remove")
def remove_from_briefing(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    title = payload.get("title", "")
    if not title:
        return {"status": "error", "message": "No title provided"}

    latest = get_latest_briefing_file_for_profile(profile)
    if not latest:
        return {"status": "error", "message": "No briefing file found"}
    try:
        with open(latest, "r", encoding="utf-8") as f:
            data = json.load(f)
        original_count = len(data)
        data = [item for item in data if item.get("title", "") != title]
        removed = original_count - len(data)
        if removed > 0:
            with open(latest, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=4, ensure_ascii=False)
            print(f"[BRIEFING] Removed '{title[:50]}' from {os.path.basename(latest)}")
        return {"status": "success", "removed": removed, "remaining": len(data), "profile": profile}
    except Exception as e:
        return {"status": "error", "message": str(e)}


@app.post("/briefing/restore")
def restore_to_briefing(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    article = payload.get("article")
    if not article:
        return {"status": "error", "message": "No article provided"}

    latest = get_latest_briefing_file_for_profile(profile)
    if not latest:
        return {"status": "error", "message": "No briefing file found"}
    try:
        with open(latest, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not any(item.get("title") == article.get("title") for item in data):
            data.insert(0, article)
            with open(latest, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=4, ensure_ascii=False)
            print(f"[BRIEFING] Restored '{article.get('title', '')[:50]}' to {os.path.basename(latest)}")
        return {"status": "success", "count": len(data), "profile": profile}
    except Exception as e:
        return {"status": "error", "message": str(e)}


# ==========================================
# --- HISTORY ENDPOINTS ---
# ==========================================
@app.get("/history/list")
def get_history_list(request: Request, session_id: str = Query(None)):
    profile = get_active_profile_name(request)
    files = get_profile_history_files(profile)
    files.sort(key=os.path.getmtime, reverse=True)
    file_list = []
    for f in files:
        filename = os.path.basename(f)
        try:
            if filename.startswith("briefing_"):
                ts = filename.replace("briefing_", "").replace(".json", "")
                display_date = datetime.datetime.strptime(ts, "%Y-%m-%d_%H-%M-%S").strftime("%b %d, %Y - %I:%M %p")
                file_list.append({"filename": filename, "display": display_date, "type": "scheduler", "profile": profile})
            elif filename.startswith("manual_"):
                if not session_id: continue
                expected_prefix = f"manual_{session_id}_"
                if not filename.startswith(expected_prefix): continue
                ts = filename.replace(expected_prefix, "").replace(".json", "")
                display_date = datetime.datetime.strptime(ts, "%Y-%m-%d_%H-%M-%S").strftime("%b %d, %Y - %I:%M %p")
                file_list.append({"filename": filename, "display": display_date, "type": "manual", "profile": profile})
        except:
            pass
    return file_list


@app.get("/history/range")
def get_history_by_range(request: Request, from_date: str, to_date: str, session_id: str = Query(None)):
    profile = get_active_profile_name(request)
    try:
        start_date = datetime.datetime.strptime(from_date, "%Y-%m-%d").date()
        end_date = datetime.datetime.strptime(to_date, "%Y-%m-%d").date()
    except ValueError:
        return {"status": "error", "message": "Invalid date format.", "profile": profile}

    merged_results = []
    seen_titles = set()
    files = get_profile_history_files(profile)
    for f in files:
        filename = os.path.basename(f)
        try:
            if filename.startswith("briefing_"):
                ts = filename.replace("briefing_", "").replace(".json", "")
                date_part = ts.split("_")[0]
            elif filename.startswith("manual_"):
                if not session_id: continue
                expected_prefix = f"manual_{session_id}_"
                if not filename.startswith(expected_prefix): continue
                ts = filename.replace(expected_prefix, "").replace(".json", "")
                date_part = ts.split("_")[0]
            else:
                continue

            file_date = datetime.datetime.strptime(date_part, "%Y-%m-%d").date()
            if start_date <= file_date <= end_date:
                with open(f, "r", encoding="utf-8") as file_obj:
                    data = json.load(file_obj)
                    for item in data:
                        item = apply_learned_region(item, profile)
                        title = item.get("title", "")
                        if title not in seen_titles:
                            seen_titles.add(title)
                            merged_results.append(item)
        except:
            continue

    merged_results.sort(key=lambda x: x.get("date", ""), reverse=True)
    return {"status": "success", "count": len(merged_results), "results": merged_results, "profile": profile}


@app.get("/history/{filename}")
def get_history_file(request: Request, filename: str):
    profile = get_active_profile_name(request)
    file_path = resolve_profile_history_file(filename, profile)
    if not file_path:
        raise HTTPException(status_code=404, detail="File not found")
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return {"status": "success", "results": apply_learned_regions(data, profile), "profile": profile}
    except Exception as e:
        return {"status": "error", "message": str(e)}


# ==========================================
# VOC FEEDBACK STORE
# ==========================================
VOC_FEEDBACK_FILE = os.path.join(ROOT_DIR, "voc_feedback.json")


@app.post("/voc")
async def submit_voc_feedback(request: Request):
    import uuid

    payload = await request.json()
    message = str(payload.get("message", "")).strip()
    if not message:
        return {"status": "error", "message": "Feedback message is required"}
    items = []
    if os.path.exists(VOC_FEEDBACK_FILE):
        try:
            with open(VOC_FEEDBACK_FILE, "r", encoding="utf-8") as f:
                items = json.load(f)
        except Exception:
            items = []
    profile = get_active_profile_name(request)
    feedback_item = {
        "id": str(uuid.uuid4()),
        "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
        "name": str(payload.get("name", "anonymous")).strip() or "anonymous",
        "type": str(payload.get("type", "ui_feedback")).strip(),
        "message": message,
        "page": str(payload.get("page", "unknown")).strip(),
        "profile": profile,
        "client_host": request.client.host if request.client else None,
    }
    items.insert(0, feedback_item)
    with open(VOC_FEEDBACK_FILE, "w", encoding="utf-8") as f:
        json.dump(items[:500], f, indent=2, ensure_ascii=False)
    return {"status": "success", "message": "Feedback saved", "item": feedback_item}


# ==========================================
# --- TRAINING ENDPOINTS ---
# ==========================================
@app.post("/train")
def save_training_data(request: Request, data: VotePayload, background_tasks: BackgroundTasks):
    profile = get_active_profile_name(request)
    total = save_training_vote(data.keywords, data.summary, data.vote, data.title or "", profile)
    background_tasks.add_task(retrain_and_reload, profile)
    return {"status": "success", "total_samples": total, "profile": profile, "retrain_scheduled": True}


# ==========================================
# --- NOT INTERESTED ENDPOINTS ---
# ==========================================
@app.post("/not-interested")
def add_not_interested(request: Request, background_tasks: BackgroundTasks, payload: dict = Body(...)):
    profile = get_active_profile_name(request)
    title = payload.get("title", "")
    summary = (
        payload.get("master_summary")
        or payload.get("summary")
        or payload.get("snippet")
        or payload.get("full_content")
        or payload.get("full_contents")
        or ""
    )
    keywords = payload.get("keywords_found", [])

    with not_interested_lock:
        store = load_not_interested_store(request)
        if is_already_rejected(title, store):
            return {"status": "exists", "message": "Already in Not Interested", "count": len(store), "profile": profile}

        entry = {
            "title": title,
            "master_summary": summary,
            "ppt_summary": payload.get("ppt_summary", ""),
            "snippet": payload.get("snippet", ""),
            "date": payload.get("date", ""),
            "link": payload.get("link", ""),
            "top_image": payload.get("top_image", ""),
            "sources": payload.get("sources", []),
            "importance_score": payload.get("importance_score", 0),
            "keywords_found": keywords,
            "region": payload.get("region", "Global"),
            "full_contents": payload.get("full_contents", ""),
            "category": payload.get("category", "Tech News"),
            "source": payload.get("source", "Unknown"),
            "sentiment": payload.get("sentiment", "neutral"),
            "is_fresh": payload.get("is_fresh", True),
            "first_seen": payload.get("first_seen", ""),
            "source_count": payload.get("source_count", 1),
            "entities": payload.get("entities", {}),
            "profile": profile,
            "rejected_at": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "rejected_by": payload.get("rejected_by", "unknown"),
        }
        store.append(entry)
        save_not_interested_store(store, request)

    save_training_vote(keywords, summary, "not_interested", title, profile)
    background_tasks.add_task(retrain_and_reload, profile)

    return {"status": "success", "message": "Moved to Not Interested", "count": len(store), "profile": profile, "retrain_scheduled": True}


@app.get("/not-interested")
def get_not_interested(request: Request):
    profile = get_active_profile_name(request)
    store = load_not_interested_store(request)
    return {"status": "success", "items": apply_learned_regions(store, profile), "count": len(store), "expiry_hours": NOT_INTERESTED_EXPIRY_HOURS, "profile": profile}


@app.post("/not-interested/restore")
def restore_from_not_interested(request: Request, background_tasks: BackgroundTasks, payload: dict = Body(...)):
    profile = get_active_profile_name(request)
    title = payload.get("title", "")

    with not_interested_lock:
        store = load_not_interested_store(request)
        article_to_restore = None
        remaining = []
        for item in store:
            if item.get("title", "").strip().lower() == title.strip().lower():
                article_to_restore = item
            else:
                remaining.append(item)

        if not article_to_restore:
            return {"status": "error", "message": "Article not found in Not Interested", "profile": profile}
        save_not_interested_store(remaining, request)

    summary = article_to_restore.get("master_summary", "")
    keywords = article_to_restore.get("keywords_found", [])
    save_training_vote(keywords, summary, "interested", article_to_restore.get("title", ""), profile)

    print(f"Restored: {title[:60]}. Counter-vote saved. Triggering retrain...")
    background_tasks.add_task(retrain_and_reload, profile)

    article_to_restore.pop("rejected_at", None)
    article_to_restore.pop("rejected_by", None)

    article_to_restore = apply_learned_region(article_to_restore, profile)
    return {"status": "success", "message": "Restored to main feed", "article": article_to_restore, "count": len(remaining), "profile": profile, "retrain_scheduled": True}


# ==========================================
# --- USAGE TRACKING ENDPOINTS ---
# ==========================================
@app.post("/track")
def track_activity(payload: dict = Body(...), request: Request = None):
    ip = get_client_ip(request)
    profile = get_active_profile_name(request)
    team_owner = get_team_owner_for_ip(ip)
    fingerprint = payload.get("fingerprint", "unknown")
    action = payload.get("action", "")
    detail = payload.get("detail", "")

    if not action:
        return {"status": "ok"}

    device_id = get_device_id(ip, fingerprint)
    today = get_today()

    with tracker_lock:
        tracker = load_tracker()

        if device_id not in tracker:
            tracker[device_id] = {
                "ip": ip,
                "fingerprint": fingerprint,
                "profile": profile,
                "owner": team_owner or "Unknown",
                "known_team_member": bool(team_owner),
                "first_seen": today,
                "last_seen": today,
                "activity": {},
            }

        device = tracker[device_id]
        device["last_seen"] = today
        device["ip"] = ip
        device["profile"] = profile
        device["owner"] = team_owner or "Unknown"
        device["known_team_member"] = bool(team_owner)

        if today not in device.get("activity", {}):
            device["activity"][today] = get_empty_day()

        day = device["activity"][today]

        if action == "page_load": day["page_loads"] += 1
        elif action == "search":
            if detail and detail not in day["searches"]:
                day["searches"].append(detail)
        elif action == "article_click": day["articles_clicked"] += 1
        elif action == "vote_interested": day["votes_interested"] += 1
        elif action == "vote_not_interested": day["votes_not_interested"] += 1
        elif action == "export":
            if detail and detail not in day["exports"]:
                day["exports"].append(detail)
        elif action == "briefing_view": day["briefing_views"] += 1
        elif action == "heartbeat": day["heartbeats"] += 1
        elif action == "voc_feedback":
            day.setdefault("voc_feedback", []).append(detail)

        purge_old_entries(device)
        save_tracker(tracker)

    return {"status": "ok"}


@app.get("/analytics/access")
def get_analytics_access(request: Request):
    ip = get_client_ip(request)
    allowed = is_analytics_allowed_ip(ip)

    return {
        "allowed": allowed,
        "ip": ip,
        "owner": get_team_owner_for_ip(ip) or "Unknown",
        "known_team_member": bool(get_team_owner_for_ip(ip)),
    }


@app.get("/analytics")
def get_analytics(request: Request, key: str = Query(None)):
    ip = require_analytics_access(request, key)

    tracker = load_tracker()
    today = get_today()
    summary = []

    for device_id, device in tracker.items():
        activity = device.get("activity", {})
        total_loads = sum(d.get("page_loads", 0) for d in activity.values())
        total_searches = sum(len(d.get("searches", [])) for d in activity.values())
        total_clicks = sum(d.get("articles_clicked", 0) for d in activity.values())
        total_votes = sum(d.get("votes_interested", 0) + d.get("votes_not_interested", 0) for d in activity.values())
        total_exports = sum(len(d.get("exports", [])) for d in activity.values())
        total_heartbeats = sum(d.get("heartbeats", 0) for d in activity.values())
        total_voc = sum(len(d.get("voc_feedback", [])) for d in activity.values())
        today_data = activity.get(today, get_empty_day())
        active_days = len([d for d in activity.values() if d.get("page_loads", 0) > 0])

        engagement = (total_clicks * 3 + total_votes * 5 + total_searches * 4 + total_exports * 10 + total_heartbeats * 1)

        summary.append({
            "device_id": device_id,
            "ip": device.get("ip", "unknown"),
            "owner": get_team_owner_for_ip(device.get("ip")) or device.get("owner", "Unknown"),
            "known_team_member": bool(get_team_owner_for_ip(device.get("ip"))),
            "profile": device.get("profile", DEFAULT_PROFILE),
            "first_seen": device.get("first_seen", ""),
            "last_seen": device.get("last_seen", ""),
            "active_days": active_days,
            "today": today_data,
            "totals": {
                "page_loads": total_loads,
                "searches": total_searches,
                "articles_clicked": total_clicks,
                "votes": total_votes,
                "exports": total_exports,
                "minutes_approx": total_heartbeats,
                "voc_feedback": total_voc,
            },
            "engagement_score": engagement,
            "daily": activity,
        })

    summary.sort(key=lambda x: x["engagement_score"])
    known_count = sum(1 for device in summary if device.get("known_team_member"))
    return {
        "status": "success",
        "device_count": len(summary),
        "known_team_member_count": known_count,
        "unknown_device_count": len(summary) - known_count,
        "date": today,
        "viewer": {
            "ip": ip,
            "owner": get_team_owner_for_ip(ip) or "Unknown",
        },
        "team_ip_map": TEAM_IP_MAP,
        "devices": summary,
    }


# ==========================================
# --- MANUAL CRAWL ENDPOINT (STREAMING) ---
# ==========================================
@app.get("/crawl")
async def crawl(
    request: Request,
    keywords: str = Query(None),
    from_date: str = Query(None),
    to_date: str = Query(None),
    target_sites: str = Query("All"),
    session_id: str = Query(None),
):
    profile = get_profile_for_request(request)
    sites_file = get_sites_file_for_profile(profile)
    job_id = secrets.token_hex(8)
    output_file = os.path.join(ROOT_DIR, f"ui_results_{job_id}.json")
    cluster_file = os.path.join(ROOT_DIR, f"clustered_results_{job_id}.json")

    with scheduler_lock:
        blocked_by_scheduler = SCHEDULER_STATUS["is_active"]
        active_jobs[job_id] = {
            "status": "blocked" if blocked_by_scheduler else "queued",
            "keywords": keywords,
            "started_at": datetime.datetime.now().isoformat(),
            "profile": profile,
        }

    def event_stream():
        if blocked_by_scheduler:
            yield f"data: {json.dumps({'type': 'error', 'message': 'The scheduled briefing is running now. Please start the deep scan again when it completes.'})}\n\n"
            return

        if not crawl_semaphore.acquire(blocking=False):
            yield f"data: {json.dumps({'type': 'error', 'message': 'Server is at capacity. Please wait a moment and try again.'})}\n\n"
            with scheduler_lock:
                active_jobs[job_id]["status"] = "error"
            return

        process = None

        try:
            with scheduler_lock:
                active_jobs[job_id]["status"] = "running"
            yield f"data: {json.dumps({'type': 'job_started', 'job_id': job_id, 'profile': profile})}\n\n"
            yield f"data: {json.dumps({'type': 'status', 'message': f'Using {profile} profile'})}\n\n"
            yield f"data: {json.dumps({'type': 'status', 'message': 'Deploying Spider...'})}\n\n"

            cmd = [
                sys.executable, "-m", "scrapy", "crawl", "news_spider",
                "-a", f"keyword={keywords}",
                "-a", f"from_date={from_date}",
                "-a", f"to_date={to_date}",
                "-a", f"target_sites={target_sites}",
                "-a", f"sites_file={sites_file}",
                "-s", "ROBOTSTXT_OBEY=False",
                "-s", "USER_AGENT=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
                "-s", "DNS_RESOLVER=scrapy.resolver.CachingThreadedResolver",
                "-s", "TWISTED_REACTOR=twisted.internet.asyncioreactor.AsyncioSelectorReactor",
                "-O", output_file,
            ]
            spider_cwd = os.path.join(ROOT_DIR, "news_aggregator")

            process = subprocess.Popen(
                cmd, cwd=spider_cwd,
                stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                text=True, bufsize=1, universal_newlines=True,
                encoding="utf-8", errors="replace",
                close_fds=CLOSE_FDS,
            )

            try:
                for line in process.stdout:
                    line = line.strip()
                    if not line: continue
                    sys.stdout.write(f"{line}\n")
                    sys.stdout.flush()
                    if "LOG:" in line:
                        yield f"data: {json.dumps({'type': 'status', 'message': line.split('LOG:', 1)[1].strip()})}\n\n"
                    elif "item_scraped_count" in line:
                        yield f"data: {json.dumps({'type': 'status', 'message': 'Gathering Intelligence...'})}\n\n"
                process.wait()
            finally:
                if process and process.stdout:
                    process.stdout.close()

            # AI Bouncer
            yield f"data: {json.dumps({'type': 'status', 'message': 'Running AI Gatekeeper...'})}\n\n"

            if os.path.exists(output_file):
                try:
                    with open(output_file, "r", encoding="utf-8") as f:
                        raw_data = json.load(f)

                    filtered_data, dropped_count, low_priority_count = run_bouncer_filter_on_items(
                        raw_data, profile, "manual_raw"
                    )

                    with open(output_file, "w", encoding="utf-8") as f:
                        json.dump(filtered_data, f, indent=4, ensure_ascii=False)

                    yield f"data: {json.dumps({'type': 'status', 'message': f'Gatekeeper done. Removed {dropped_count} articles. Low priority kept: {low_priority_count}.'})}\n\n"
                    print(f"Bouncer complete [{profile}]. Dropped {dropped_count} articles.", flush=True)
                except Exception as e:
                    print(f"Bouncer error, skipping filter: {e}", flush=True)

            # ==========================================
            # PHASE 1: STREAM CARDS IN REAL-TIME
            # ==========================================
            yield f"data: {json.dumps({'type': 'status', 'message': 'Activating Fusion Engine (Streaming Mode)...'})}\n\n"

            try:
                from semantic_clustering import MinimalSemanticEngine
                engine = MinimalSemanticEngine()

                streamed_count = 0
                for event in engine.fuse_stream(job_id=job_id, fast_mode=False):
                    event["category"] = assign_category(
                        event.get("title", ""),
                        event.get("master_summary", "") or event.get("snippet", ""),
                    )
                    event["profile"] = profile
                    event.update(apply_learned_region(event, profile))
                    yield f"data: {json.dumps({'type': 'card', 'card': event}, ensure_ascii=False)}\n\n"
                    streamed_count += 1

                if streamed_count == 0:
                    yield f"data: {json.dumps({'type': 'status', 'message': 'No articles to process.'})}\n\n"
                else:
                    yield f"data: {json.dumps({'type': 'status', 'message': f'All {streamed_count} cards streamed. Optimizing...'})}\n\n"

                # PHASE 2: RE-CLUSTER
                if streamed_count > 1:
                    yield f"data: {json.dumps({'type': 'status', 'message': 'Clustering duplicate stories...'})}\n\n"
                    engine.fuse_cluster(job_id=job_id, fast_mode=False)
                    yield f"data: {json.dumps({'type': 'status', 'message': 'Optimization complete.'})}\n\n"
                else:
                    if os.path.exists(output_file):
                        with open(output_file, "r", encoding="utf-8") as f:
                            raw = json.load(f)
                        with open(cluster_file, "w", encoding="utf-8") as f:
                            json.dump(raw, f, indent=4, ensure_ascii=False)

            except Exception as e:
                print(f"Fusion streaming error: {e}", flush=True)
                yield f"data: {json.dumps({'type': 'status', 'message': 'Fusion Engine Error. Falling back...'})}\n\n"
                try:
                    fusion_fallback = subprocess.Popen(
                        [sys.executable, "-u", "semantic_clustering.py", "--job-id", job_id],
                        cwd=ROOT_DIR,
                        stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                        text=True, encoding="utf-8", errors="replace",
                        close_fds=CLOSE_FDS,
                    )
                    try:
                        for line in fusion_fallback.stdout:
                            line = line.strip()
                            sys.stdout.write(f"{line}\n")
                            sys.stdout.flush()
                            if "FUSION ENGINE:" in line:
                                yield f"data: {json.dumps({'type': 'status', 'message': line.split('FUSION ENGINE:', 1)[1].strip()})}\n\n"
                        fusion_fallback.wait()
                    finally:
                        if fusion_fallback and fusion_fallback.stdout:
                            fusion_fallback.stdout.close()
                except Exception as fallback_err:
                    print(f"Fallback fusion also failed: {fallback_err}", flush=True)

            # Load final results
            results = []
            if os.path.exists(cluster_file):
                with open(cluster_file, "r", encoding="utf-8") as f:
                    results = json.load(f)
            elif os.path.exists(output_file):
                with open(output_file, "r", encoding="utf-8") as f:
                    results = json.load(f)

            if results:
                results, final_dropped_count, final_low_priority_count = run_bouncer_filter_on_items(
                    results, profile, "manual_final"
                )
                for r in results:
                    r["category"] = assign_category(
                        r.get("title", ""),
                        r.get("master_summary", "") or r.get("snippet", ""),
                    )
                    r["profile"] = profile
                    r.update(apply_learned_region(r, profile))

            # Archive
            if results:
                yield f"data: {json.dumps({'type': 'status', 'message': 'Archiving Intelligence...'})}\n\n"
                learner.log_search_data(keywords, results)
                timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
                sid = session_id or "unknown"
                manual_path = os.path.join(get_profile_history_dir(profile), f"manual_{sid}_{timestamp}.json")
                with open(manual_path, "w", encoding="utf-8") as f:
                    json.dump(results, f, indent=4, ensure_ascii=False)
                print(f"Archived [{profile}]: manual_{sid}_{timestamp}.json", flush=True)

            with scheduler_lock:
                active_jobs[job_id]["status"] = "complete"
            yield f"data: {json.dumps({'type': 'data', 'results': results, 'job_id': job_id, 'reclustered': True, 'profile': profile})}\n\n"

        finally:
            if process and process.poll() is None:
                process.terminate()
            crawl_semaphore.release()
            with scheduler_lock:
                active_jobs[job_id]["status"] = "complete"
            threading.Timer(300, cleanup_job_files, args=[output_file, cluster_file]).start()

    return StreamingResponse(event_stream(), media_type="text/event-stream")


# ==========================================
# --- FRONTEND ROUTING ---
# ==========================================
API_ROUTES = {
    "crawl", "train", "status", "briefing",
    "export-excel", "export-ppt", "export-word",
    "sites", "latest-briefing", "workflow",
    "history", "not-interested", "region", "track", "analytics", "profile", "voc",
    "insight",
}


@app.get("/")
def serve_root():
    index_path = os.path.join(abs_frontend_path, "index.html")
    if os.path.exists(index_path):
        return FileResponse(index_path)
    return {"status": "error", "message": "UI not built yet. Check frontend/dist folder."}


@app.get("/{catchall:path}")
def serve_react_app(catchall: str):
    root = catchall.split("/")[0]
    if root in API_ROUTES:
        raise HTTPException(status_code=404, detail="Not Found")
    index_path = os.path.join(abs_frontend_path, "index.html")
    if os.path.exists(index_path):
        return FileResponse(index_path)
    return {"status": "error", "message": "UI not built yet. Check frontend/dist folder."}
